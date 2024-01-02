import argparse
import os

import PyPDF2
import docx
import fitz
import httpx
import openpyxl
import pptx
from lxml import etree
from openai import OpenAI
from pptx import Presentation


class DocumentProcessor:
    def __init__(self, file_path):
        self.file_path = file_path
        self.file_extension = os.path.splitext(file_path)[-1].lower()
        self.file_name = os.path.splitext(file_path)[0]

    def extract_text(self):
        if self.file_extension == '.docx':
            return self._extract_from_docx()
        elif self.file_extension == '.txt' or self.file_extension == '.csv' or self.file_extension == '.json':
            return self._extract_from_text_file()
        elif self.file_extension == '.pdf':
            return self._extract_from_pdf()
        elif self.file_extension == '.xlsx':
            return self._extract_from_xlsx()
        elif self.file_extension == '.pptx':
            return self._extract_from_pptx()
        else:
            raise ValueError(f"Unsupported file type: {self.file_extension}")

    def _extract_from_pptx(self):
        presentation = pptx.Presentation(self.file_path)
        full_text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return self._remove_blank_lines("\n".join(full_text))

    def _extract_from_docx(self):
        doc = docx.Document(self.file_path)
        full_text = [para.text for para in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    full_text.append(cell.text)
        return self._remove_blank_lines("\n".join(full_text))

    def _extract_from_text_file(self):
        with open(self.file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        return self._remove_blank_lines(content)

    def _extract_from_pdf(self):
        with fitz.open(self.file_path) as doc:
            text = ""
            for page in doc:
                text += page.get_text()
        return self._remove_blank_lines(text)

    def _extract_from_xlsx(self):
        workbook = openpyxl.load_workbook(self.file_path)
        sheet = workbook.active
        content = "\n".join(cell[0].value for cell in sheet.iter_rows(values_only=True) if cell[0].value)
        return self._remove_blank_lines(content)

    @staticmethod
    def _remove_blank_lines(text):
        lines = [line.strip() for line in text.splitlines()]
        non_blank_lines = [line for line in lines if line]
        return '\n'.join(non_blank_lines)

    def edit_metadata(self, new_keywords):
        if self.file_extension == '.docx':
            doc = docx.Document(self.file_path)
            core_props = doc.core_properties
            core_props.keywords = new_keywords
            doc.save(self.file_path)
        if self.file_extension == '.pdf':
            reader = PyPDF2.PdfReader(self.file_path)
            writer = PyPDF2.PdfWriter()
            writer.append_pages_from_reader(reader)
            # metadata = reader.getDocumentInfo()
            new_metadata = {"/Keywords": new_keywords}
            writer.add_metadata(new_metadata)
            with open(self.file_path, 'wb') as new_file:
                writer.write(new_file)
        if self.file_extension == '.pptx':
            prs = Presentation(self.file_path)
            # 获取核心属性part
            core_props = prs.core_properties._element
            nsmap = core_props.nsmap

            # 修改关键词（标签）
            keywords_element = core_props.find('cp:keywords', nsmap)
            if keywords_element is not None:
                keywords_element.text = new_keywords
            else:
                etree.SubElement(core_props,
                                 '{http://schemas.openxmlformats.org/package/2006/metadata/core-properties}keywords').text = new_keywords

            # 保存修改后的文件
            prs.save(self.file_path)

        else:
            print("Metadata editing is not supported for this file type.")


def get_openai_labels(api_key, proxy, text):
    client = OpenAI(api_key=api_key, http_client=httpx.Client(proxy=proxy))
    try:
        chat_completion = client.chat.completions.create(
            seed=549,
            messages=[
                {"role": "system",
                 "content": "You are an expert at summarizing article tags. Based on the input text, you will "
                            "summarize and return 10 best tags with independent content. Tags are separated by \",\" "
                            "and there is no number before the labels."},
                {
                    "role": "user",
                    "content": f"文本：\"{text}\". 10个最有概括力的没有数字前缀的标签：",
                }
            ],
            model="gpt-3.5-turbo",
        )
        return chat_completion.choices[0].message.content
    except Exception as e:
        print(f"Error while getting labels from OpenAI: {e}")
        return None


def main(docx_path, api_key, proxy):
    # TODO: Replace with your actual paths, keys and proxies

    if proxy is None:
        proxy = 'socks5://127.0.0.1:7890'

    processor = DocumentProcessor(docx_path)
    docx_text = processor.file_name + " :" + processor.extract_text()
    text_to_classify = docx_text[:500]
    print('提取的前500字：', text_to_classify)

    label_str = get_openai_labels(api_key, proxy, text_to_classify)
    print(type(label_str), label_str)
    if label_str:
        tags = label_str.split("、")
        keywords = ", ".join(tags)
        processor.edit_metadata(keywords)
        print(f"Document updated with keywords: {keywords}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Document Labeling Script")

    # 设置默认文件路径和API密钥
    default_file_path = 'TXT/一文搞懂华为IPD：没有流程支撑的策略，都是纸上谈兵.pdf'
    with open('openaiAPI.txt', 'r') as apifile:
        default_api_key = apifile.read()
        print('api_key:', default_api_key)

    # 添加可选的文件路径和API密钥参数
    parser.add_argument("--file_path", default=default_file_path, help="Path to the DOCX file to label")
    parser.add_argument("--api_key", default=default_api_key, help="OpenAI API Key")
    parser.add_argument("--proxy", default=None, help="Proxy setting if required")

    args = parser.parse_args()

    main(args.file_path, args.api_key, args.proxy)
